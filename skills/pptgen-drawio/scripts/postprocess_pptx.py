#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""导出后修补 .pptx：修中文不换行的坑，顺便把讲稿写进备注区。

用法：
    python scripts/postprocess_pptx.py out.pptx
    python scripts/postprocess_pptx.py out.pptx --notes paper-defense-script.md

## 为什么必须跑这一步

drawio2pptx 0.0.7 `io/drawio_loader.py` 里有一段「nowrap 启发式」：
如果一个文本框里所有段落都**不含空白字符**，它就把 `whiteSpace=wrap`
覆盖成 `wrap="none"`。

中文句子没有空格，所以**每一个纯中文文本框都会被判成不换行**。
后果是文字沿一行无限向右流，压穿右边的栏、跑出画布。
字号小的时候一行刚好放得下，看不出来；字号一大就全崩。

这个脚本把所有文本框的 `word_wrap` 强制设回 True。
`render_check.py` 的扫描是按「会换行」算的，所以不跑这一步，
扫描结果和实际渲染不是一回事。

## 讲稿写进备注

`--notes xxx.md` 按 `## 第 N 页` 切分，写到对应页的备注区。
Draw.io 没有备注这个概念，只能导出后补。
"""

import argparse
import io
import os
import re
import sys

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except AttributeError:
    pass

try:
    from pptx import Presentation
except ImportError:
    sys.exit("需要 python-pptx：pip install python-pptx")


def fix_wrap(prs):
    """强制所有文本框换行。返回改了多少个框。"""
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.word_wrap = True
                n += 1
    return n


def parse_notes(path):
    """从讲稿 md 里按 `## 第 N 页` 抽出每页要说的话，返回 {页号: 文本}。"""
    text = io.open(path, encoding="utf-8").read()
    notes, cur, buf = {}, None, []
    for line in text.splitlines():
        m = re.match(r"^##\s*第\s*(\d+)\s*页", line)
        if m:
            if cur is not None:
                notes[cur] = "\n".join(buf).strip()
            cur, buf = int(m.group(1)), []
            continue
        if cur is not None:
            buf.append(line)
    if cur is not None:
        notes[cur] = "\n".join(buf).strip()
    return notes


def apply_notes(prs, notes):
    """把讲稿写进备注区。返回写了多少页。"""
    n = 0
    for idx, slide in enumerate(prs.slides, 1):
        if idx in notes and notes[idx]:
            slide.notes_slide.notes_text_frame.text = notes[idx]
            n += 1
    return n


def main():
    ap = argparse.ArgumentParser(description="导出后修补 pptx")
    ap.add_argument("pptx")
    ap.add_argument("--notes", help="讲稿 md，按「## 第 N 页」写进备注区")
    ap.add_argument("-o", "--out", help="输出文件，默认原地覆盖")
    args = ap.parse_args()

    if not os.path.exists(args.pptx):
        sys.exit("找不到文件：%s" % args.pptx)

    prs = Presentation(args.pptx)
    n_wrap = fix_wrap(prs)
    print("✅ %d 个文本框强制换行（修 drawio2pptx 中文 nowrap 的坑）" % n_wrap)

    if args.notes:
        if not os.path.exists(args.notes):
            sys.exit("找不到讲稿：%s" % args.notes)
        n_notes = apply_notes(prs, parse_notes(args.notes))
        print("✅ %d 页写入备注" % n_notes)

    out = args.out or args.pptx
    prs.save(out)
    print("已保存 %s" % out)
    print("接下来跑 render_check.py 看图确认。")


if __name__ == "__main__":
    main()
