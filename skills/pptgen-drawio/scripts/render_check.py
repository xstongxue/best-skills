#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""把 .pptx 每页导成 PNG，并扫描文字溢出。

用法：
    python scripts/render_check.py out.pptx
    python scripts/render_check.py out.pptx --outdir render --dpi 110

做两件事：

1. **溢出扫描**（自动）——用 Pillow 加载真实字体文件量每一行的像素宽度，
   算出文本框实际需要多高、多宽，和框的尺寸比。比「行高约 1.4 倍字号」这种
   估算准，因为中文和英文字宽不一样，微软雅黑和宋体也不一样。
2. **导出 PNG**（供人和 Agent 看）——Windows 上走 PowerPoint COM，
   其他平台走 LibreOffice + pdftoppm。哪个都没有就只做第 1 步。

退出码：有 OVERFLOW 级别的问题返回 1，否则 0。
"""

import argparse
import os
import platform
import shutil
import subprocess
import sys

# Windows 控制台默认 GBK，打不出 ⚠️ / ❌ 之类的字符，先把 stdout 换成 UTF-8
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except AttributeError:      # Python < 3.7
    import codecs
    sys.stdout = codecs.getwriter("utf-8")(sys.stdout.buffer, "replace")
    sys.stderr = codecs.getwriter("utf-8")(sys.stderr.buffer, "replace")

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("需要 python-pptx：pip install python-pptx")

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from textmetrics import (          # noqa: E402  量文字用的公共模块
    PAD_X_PT, find_font_path, load_font, missing_fonts,
    measure_pt, text_width_pt, wrap_lines,
)


# ---------------------------------------------------------------- 溢出扫描

# 连续空白超过这么多 pt（约 240 px）就提醒排版没铺开
SPARSE_GAP_PT = 180.0


def _emu_to_pt(v):
    return Emu(v).pt if v is not None else None


def _largest_gap(bands, slide_h_pt):
    """把所有元素的纵向区间合并，找出中间最大的一段空白。
    返回 (起始 y, 高度)，单位 pt。只看元素之间的空，不算页面上下边缘。"""
    if not bands:
        return 0.0, 0.0
    merged = []
    for lo, hi in sorted(bands):
        if merged and lo <= merged[-1][1]:
            merged[-1][1] = max(merged[-1][1], hi)
        else:
            merged.append([lo, hi])
    best_top, best_h = 0.0, 0.0
    for i in range(len(merged) - 1):
        gap = merged[i + 1][0] - merged[i][1]
        if gap > best_h:
            best_top, best_h = merged[i][1], gap
    return best_top, best_h


def scan_overflow(pptx_path, strict_ratio=1.0, warn_ratio=0.9):
    """逐页逐文本框量一遍。返回 (问题列表, 统计)。"""
    prs = Presentation(pptx_path)
    slide_w_pt = Emu(prs.slide_width).pt
    slide_h_pt = Emu(prs.slide_height).pt
    issues = []
    n_boxes = 0

    for idx, slide in enumerate(prs.slides, 1):
        bands = []          # 这一页所有非满版元素占的纵向区间，用来找大片空白
        n_text = 0
        for shape in slide.shapes:
            box_w = _emu_to_pt(shape.width) or 0
            box_h = _emu_to_pt(shape.height) or 0
            left = _emu_to_pt(shape.left) or 0
            top = _emu_to_pt(shape.top) or 0

            # 满版背景块不算内容
            if box_w * box_h < 0.85 * slide_w_pt * slide_h_pt:
                bands.append((top, top + box_h))

            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if not text.strip():
                continue
            n_boxes += 1
            n_text += 1

            # 取这个框里最大的字号和第一个字体名
            size_pt, fname = 18.0, "微软雅黑"
            sizes = []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        sizes.append(run.font.size.pt)
                    if run.font.name:
                        fname = run.font.name
            if sizes:
                size_pt = max(sizes)

            font = load_font(fname, size_pt)
            # 内边距读实际值——drawio2pptx 写的是 0，按 python-pptx 默认值
            # 算会多算 7 pt，单行文字全被误报成溢出。
            tf = shape.text_frame
            pad_x = sum(_emu_to_pt(v) or 0
                        for v in (tf.margin_left, tf.margin_right))
            pad_y = sum(_emu_to_pt(v) or 0
                        for v in (tf.margin_top, tf.margin_bottom))
            avail_w = box_w - pad_x
            lines, need_h = measure_pt(text, size_pt, box_w, fname,
                                       pad_x_pt=pad_x, pad_y_pt=pad_y)

            preview = text.replace("\n", "⏎")[:34]
            ratio_h = need_h / box_h if box_h else 99

            if ratio_h > strict_ratio:
                issues.append(dict(
                    level="OVERFLOW", slide=idx, kind="高度不够",
                    detail="需 %.0fpt / 框高 %.0fpt（%d 行 × %.0fpt）"
                           % (need_h, box_h, lines, size_pt),
                    ratio=ratio_h, font=fname, text=preview))
            elif ratio_h > warn_ratio and lines > 1:
                # 单行文字刚好填满一行高的框是正常的，不报。
                # 只有多行文字快顶满才值得提醒——再多一行就裁字了。
                issues.append(dict(
                    level="TIGHT", slide=idx, kind="几乎顶满",
                    detail="需 %.0fpt / 框高 %.0fpt（%d 行，再多一行就溢出）"
                           % (need_h, box_h, lines),
                    ratio=ratio_h, font=fname, text=preview))

            # 单行放不下且框内不折行的情况：宽度就是不够
            if lines == 1 and text_width_pt(text, font, size_pt) > avail_w:
                issues.append(dict(
                    level="OVERFLOW", slide=idx, kind="宽度不够",
                    detail="需 %.0fpt / 可用 %.0fpt"
                           % (text_width_pt(text, font, size_pt), avail_w),
                    ratio=0, font=fname, text=preview))

            # 出画布
            if left < -1 or top < -1 or left + box_w > slide_w_pt + 1 \
                    or top + box_h > slide_h_pt + 1:
                issues.append(dict(
                    level="OVERFLOW", slide=idx, kind="超出画布",
                    detail="框 (%.0f,%.0f)+(%.0f×%.0f) vs 画布 %.0f×%.0f"
                           % (left, top, box_w, box_h, slide_w_pt, slide_h_pt),
                    ratio=0, font=fname, text=preview))

        # 大片空白检查。节标题页本来就空，只查文本框 ≥6 个的内容页。
        if n_text >= 6:
            gap_top, gap_h = _largest_gap(bands, slide_h_pt)
            if gap_h > SPARSE_GAP_PT:
                issues.append(dict(
                    level="SPARSE", slide=idx, kind="大片空白",
                    detail="y %.0f–%.0fpt 连续空 %.0fpt（约 %.0f px），"
                           "内容没铺开或下半页是空的"
                           % (gap_top, gap_top + gap_h, gap_h, gap_h / 0.75),
                    ratio=0, font="-", text=""))

    stats = dict(slides=len(prs.slides._sldIdLst),
                 boxes=n_boxes,
                 missing_fonts=missing_fonts())
    return issues, stats


# ---------------------------------------------------------------- 导出 PNG

def export_powerpoint_com(pptx_path, outdir, width=1600):
    """Windows：用 PowerPoint COM 导出每页 PNG。"""
    import win32com.client
    pptx_path = os.path.abspath(pptx_path)
    outdir = os.path.abspath(outdir)
    app = None
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        # 有些版本不允许把 Visible 设为 False，忽略即可
        pres = app.Presentations.Open(pptx_path, WithWindow=False)
        try:
            h = int(width * pres.PageSetup.SlideHeight / pres.PageSetup.SlideWidth)
            pres.Export(outdir, "PNG", width, h)
            n = pres.Slides.Count
        finally:
            pres.Close()
        return n
    finally:
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass


def export_soffice(pptx_path, outdir, dpi=110):
    """跨平台：LibreOffice 转 PDF，再 pdftoppm 切图。"""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice:
        raise RuntimeError("找不到 soffice")
    subprocess.run([soffice, "--headless", "--convert-to", "pdf",
                    "--outdir", outdir, pptx_path], check=True,
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    pdf = os.path.join(outdir, os.path.splitext(
        os.path.basename(pptx_path))[0] + ".pdf")
    if not pdftoppm:
        raise RuntimeError("转出了 PDF 但找不到 pdftoppm：%s" % pdf)
    subprocess.run([pdftoppm, "-jpeg", "-r", str(dpi), pdf,
                    os.path.join(outdir, "page")], check=True)
    return len([f for f in os.listdir(outdir)
                if f.startswith("page") and f.endswith(".jpg")])


def render(pptx_path, outdir, dpi=110, width=1600):
    os.makedirs(outdir, exist_ok=True)
    if platform.system() == "Windows":
        try:
            n = export_powerpoint_com(pptx_path, outdir, width)
            return "PowerPoint COM", n
        except Exception as e:
            print("  PowerPoint COM 导出失败：%s" % e)
    try:
        n = export_soffice(pptx_path, outdir, dpi)
        return "LibreOffice + pdftoppm", n
    except Exception as e:
        print("  LibreOffice 导出失败：%s" % e)
    return None, 0


# ---------------------------------------------------------------- main

def main():
    ap = argparse.ArgumentParser(description="pptx 溢出扫描 + 逐页导图")
    ap.add_argument("pptx")
    ap.add_argument("--outdir", default=None, help="PNG 输出目录，默认 <pptx名>_render")
    ap.add_argument("--dpi", type=int, default=110)
    ap.add_argument("--width", type=int, default=1600, help="COM 导出宽度 px")
    ap.add_argument("--no-render", action="store_true", help="只扫描不导图")
    args = ap.parse_args()

    if not os.path.exists(args.pptx):
        sys.exit("找不到文件：%s" % args.pptx)

    print("=" * 66)
    print("溢出扫描：%s" % args.pptx)
    print("=" * 66)
    issues, stats = scan_overflow(args.pptx)

    bad = [i for i in issues if i["level"] == "OVERFLOW"]
    tight = [i for i in issues if i["level"] == "TIGHT"]
    sparse = [i for i in issues if i["level"] == "SPARSE"]

    if stats["missing_fonts"]:
        print("⚠️  这些字体本机没装，量出来的宽度只是估算：%s"
              % "、".join(stats["missing_fonts"]))

    order = {"OVERFLOW": 0, "TIGHT": 1, "SPARSE": 2}
    marks = {"OVERFLOW": "❌", "TIGHT": "⚠️ ", "SPARSE": "○ "}
    for i in sorted(issues, key=lambda x: (x["slide"], order[x["level"]])):
        print("%s 第 %2d 页 [%s] %s"
              % (marks[i["level"]], i["slide"], i["kind"], i["detail"]))
        if i["text"]:
            print("     「%s」%s" % (i["text"], i["font"]))

    print("-" * 66)
    print("%d 页 / %d 个文本框：%d 处溢出，%d 处几乎顶满，%d 页留白过大"
          % (stats["slides"], stats["boxes"], len(bad), len(tight), len(sparse)))
    if sparse:
        print("○ 留白过大不算错，但要确认是有意留白还是内容没铺开——看图确认。")

    if not args.no_render:
        outdir = args.outdir or os.path.splitext(args.pptx)[0] + "_render"
        print("\n导出 PNG 到 %s/ ..." % outdir)
        engine, n = render(args.pptx, outdir, args.dpi, args.width)
        if engine:
            print("✅ %s 导出 %d 页。**接下来逐页看图**，先看文字有没有跑出框。"
                  % (engine, n))
        else:
            print("❌ 没有可用的导出工具。装 LibreOffice 或在 Windows 上装 "
                  "PowerPoint + pywin32，否则只能靠上面的扫描结果。")

    return 1 if bad else 0


if __name__ == "__main__":
    sys.exit(main())
